"""
Generate thesis PowerPoint: Stop Relocation Scenarios (R_*).
Run: python create_relocation_presentation.py
Output: Stop_Relocation_Scenarios_Presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Stop_Relocation_Scenarios_Presentation.pptx"

# Theme colors
NAVY = RGBColor(0x1A, 0x3A, 0x5C)
TEAL = RGBColor(0x00, 0x7A, 0x87)
ACCENT = RGBColor(0xF4, 0xA2, 0x61)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(9), Inches(0.35))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = ACCENT


def add_bullets(slide, items: list[str], left=0.6, top=1.35, width=8.8, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    col_widths: list[float] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, title)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    left, top = Inches(0.45), Inches(1.35)
    width = Inches(9.1)
    height = Inches(0.45 * n_rows)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF4)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_title: str,
    left_items: list[str],
    right_title: str,
    right_items: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, title)

    for side, stitle, sitems, x in (
        ("L", left_title, left_items, 0.45),
        ("R", right_title, right_items, 5.05),
    ):
        tb = slide.shapes.add_textbox(Inches(x), Inches(1.25), Inches(4.4), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEAL

        box = slide.shapes.add_textbox(Inches(x), Inches(1.65), Inches(4.4), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(sitems):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = item
            para.font.size = Pt(14)
            para.font.color.rgb = DARK
            para.space_after = Pt(8)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(8.6), Inches(1.2))
    p = tb.text_frame.paragraphs[0]
    p.text = "Stop Relocation Scenarios"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(8.6), Inches(1.0))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Formal-Stop Placement Sensitivity Along the\nRobinsons–Waltermart Bus Corridor"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER

    tb3 = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(8.6), Inches(0.8))
    p3 = tb3.text_frame.paragraphs[0]
    p3.text = "Thesis Simulation Study  |  UXsim Microsimulation  |  Dec 2025 Field Calibration"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p3.alignment = PP_ALIGN.CENTER

    # --- Slide 2: Context ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Study Context", "Where relocation fits in the thesis")
    add_bullets(
        slide,
        [
            "Primary thesis comparison: Scenario A (Baseline) vs Scenario D (Optimized)",
            "  • A: formal + informal stops, 100% field dwell",
            "  • D: formal stops only, 55% dwell — stop-policy reform",
            "",
            "Supplementary analysis: 31 stop-relocation scenarios (R_*)",
            "  • Tests whether ±5 m formal-stop shifts change corridor travel time",
            "  • Distinct question from informal-stop removal and dwell reform",
            "",
            "63,612 relocation runs vs 4,104 primary policy runs (A + D)",
        ],
    )

    # --- Slide 3: Research question ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Research Question")
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "If formal stops remain mandatory but are shifted slightly "
        "along the road (+5 m), does corridor bus performance change meaningfully?"
    )
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    add_bullets(
        slide,
        [
            "Formal-only network (same stop compliance as Scenario C)",
            "100% field-calibrated dwell (no 0.55 scaling)",
            "+5 m shift along corridor chainage (toward Waltermart)",
            "Same 513 field buses, 4 signal patterns, per-bus traffic & crowding",
        ],
        top=4.2,
        size=16,
    )

    # --- Slide 4: Eligible stops ---
    add_table_slide(
        prs,
        "Eligible Formal Stops (S1–S5)",
        ["Code", "Stop Key", "Corridor Role"],
        [
            ["S1", "Robinsons terminal", "Origin terminal (Waltermart → Robinsons)"],
            ["S2", "RKP Trading", "Formal mid-route (Robinsons → Waltermart)"],
            ["S3", "Villa Verde", "Formal mid-route (Robinsons → Waltermart)"],
            ["S4", "Vista Mall", "Formal mid-route (Robinsons → Waltermart)"],
            ["S5", "Waltermart terminal", "Destination terminal (both directions)"],
        ],
        col_widths=[0.9, 2.2, 6.0],
    )

    # --- Slide 5: Mask encoding ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "5-Bit Relocation Mask", "Scenario IDs: R_00001 … R_11111")
    add_bullets(
        slide,
        [
            "Each bit = one stop shifted +5 m (1 = shifted, 0 = original position)",
            "",
            "Bit order (left → right):  S1  S2  S3  S4  S5",
            "                           Robinsons | RKP | Villa Verde | Vista Mall | Waltermart",
            "",
            "Examples:",
            "  • relocate_mask = 00100  →  Villa Verde (+5 m) only  →  R_00100",
            "  • relocate_mask = 10001  →  Robinsons + Waltermart shifted  →  R_10001",
            "  • relocate_mask = 11111  →  all five stops shifted  →  R_11111",
            "",
            "2⁵ − 1 = 31 non-zero combinations (no shift = Scenario C at original positions)",
        ],
        size=17,
    )

    # --- Slide 6: Scenario rules ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Rules for Every Relocation Scenario")
    items = [
        "1. Formal-only stop network — informal stops OFF (Scenario C compliance)",
        "2. Dwell at 100% of field-calibrated values (no optimized 0.55 scaling)",
        "3. Selected formal stops shifted +5 m along Robinsons-origin chainage",
        "4. When Waltermart terminal (S5) is shifted, corridor length extends so links stay consistent",
        "5. Same per-bus calibration: traffic (Light/Moderate/Heavy), crowding, boarding/alighting",
        "6. Four signal encounter patterns tested: G-G, G-R, R-G, R-R",
    ]
    add_bullets(slide, items, size=17)

    # --- Slide 7: A-D vs R comparison ---
    add_table_slide(
        prs,
        "Relocation vs Policies A–D",
        ["Feature", "Policies A–D", "Relocation R_*"],
        [
            ["Purpose", "Stop policy (who stops where; dwell rules)", "Stop position sensitivity (+5 m)"],
            ["Informal stops", "Varies by scenario", "Always OFF"],
            ["Dwell scale", "100% or 55% (Scenario D only)", "Always 100%"],
            ["Primary thesis comparison", "Yes — A vs D", "No — supplementary / Section 4.9"],
            ["Run count (full scale)", "4,104 (A + D) or 8,208 (A–D)", "63,612 (513 × 31 × 4 signals)"],
        ],
        col_widths=[2.0, 3.5, 3.6],
    )

    # --- Slide 8: Experimental design ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Experimental Design")
    add_bullets(
        slide,
        [
            "Corridor: Robinsons ↔ Waltermart, 3,255 m (extends when S5 shifted)",
            "Observations: 513 field buses from Dec 2025 data collection workbook",
            "Sessions: Morning, Lunch, Afternoon — per-bus traffic & crowding",
            "Scenarios: 31 relocation masks × 4 signal patterns = 124 runs per bus",
            "Total: 513 × 124 = 63,612 simulated corridor trips",
            "",
            "Outputs per run:",
            "  scenario_id (e.g. R_00100), scenario_group = relocation,",
            "  relocate_mask (5-bit string), sim_corridor_travel_min, sim_delay_s",
        ],
        size=16,
    )

    # --- Slide 9: What scenarios answer ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "What Relocation Scenarios Answer")
    add_bullets(
        slide,
        [
            "Do ±5 m adjustments to formal-stop chainage change travel time or delay?",
            "Does it matter which stop is shifted — terminal (S1/S5) vs mid-route bay (S2–S4)?",
            "Is placement sensitivity different for Robinsons-origin vs Waltermart-origin trips?",
            "Is fine-tuning bay coordinates secondary to policy reform (A vs D)?",
            "",
            "Expected thesis narrative:",
            "  Stop-policy change (remove informals + shorter dwell) >> micro-placement (+5 m)",
        ],
        size=17,
    )

    # --- Slide 10: Results ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Preliminary Results (63,612 Runs)", "Extended analysis — not primary thesis comparison")
    add_bullets(
        slide,
        [
            "Median corridor travel time: 16.8 min across all relocation masks",
            "All 31 masks (R_00001 … R_11111): identical median at current calibration",
            "+5 m shift per stop: ~0.00 min average effect (negligible vs traffic/signals)",
            "",
            "For comparison (stop-detail paired analysis, n = 513 buses):",
            "  Baseline → Optimized (A vs D): ~9.9 min median savings",
            "",
            "Interpretation:",
            "  Relocation is a robustness / sensitivity check — not a deployment recommendation",
            "  Thesis focus remains on stop-policy and dwell reform, not bay coordinates",
        ],
        size=16,
    )

    # --- Slide 11: Key bottlenecks (context) ---
    add_two_column_slide(
        prs,
        "Where Time Is Actually Lost (Context)",
        "High-impact factors",
        [
            "Traffic tier: Light ~14 min vs Heavy ~19 min",
            "Crowding: High +2.4 min vs Low/Medium",
            "Signal pattern: R-R +0.9 min vs G-G",
            "Informal stops + long dwell (Baseline A)",
            "RKP Trading queue (~225 s leg)",
        ],
        "Low-impact in relocation sweep",
        [
            "Shifting Robinsons +5 m (S1)",
            "Shifting RKP / Villa Verde / Vista Mall",
            "Shifting Waltermart +5 m (S5)",
            "Any combination of 1–5 stops shifted",
            "→ Supports policy-over-placement thesis",
        ],
    )

    # --- Slide 12: Section 4.9 placement ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Reporting in the Thesis (Section 4.9)")
    add_bullets(
        slide,
        [
            "Primary results (Chapter 4): Scenario A vs D — 4,104 runs",
            "  Tables: travel time, delay, speed, scenario comparison by session/traffic",
            "",
            "Extended / future work: Relocation R_* — 63,612 runs",
            "  Report mask encoding, experimental rules, and null/small effect finding",
            "  State clearly: does not replace Baseline vs Optimized comparison",
            "",
            "Suggested slide/table for defense:",
            "  One summary table: A vs D savings (~10 min) vs R_* effect (~0 min)",
        ],
        size=17,
    )

    # --- Slide 13: Conclusions ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_header_bar(slide, "Conclusions")
    add_bullets(
        slide,
        [
            "31 relocation scenarios test formal-stop placement under a formal-only rule.",
            "Design isolates +5 m chainage shifts from policy and dwell changes.",
            "Full sweep (63,612 runs) shows no meaningful travel-time difference across masks.",
            "Stop-policy reform (A → D) delivers far larger gains than bay micro-adjustment.",
            "Recommendation: prioritize informal-stop removal and dwell optimization over relocating bays by ±5 m.",
        ],
        size=18,
    )

    # --- Slide 14: Thank you ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(8.6), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(8.6), Inches(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = "Questions & Discussion"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT
    p2.alignment = PP_ALIGN.CENTER

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
